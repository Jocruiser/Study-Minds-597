from flask import Flask, request, render_template, jsonify, send_file, redirect, url_for, session, send_from_directory
from flask_cors import CORS
import os
import csv
import PyPDF2
from pptx import Presentation
import json
import requests
from dotenv import load_dotenv

app = Flask(__name__)
CORS(app)

# OpenRouter API credentials
api_url = "https://openrouter.ai/api/v1/chat/completions"
load_dotenv()  # load from .env file
api_key = os.getenv("OPENROUTER_API_KEY")
app.secret_key = os.getenv("SECRET_KEY")

headers = {
    "Authorization": f"Bearer {api_key}",
    "Content-Type": "application/json",
    "HTTP-Referer": "localhost",
    "X-Title": "Flashcard Generator",  
}

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    text_content = []
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content.append(text.strip())
    return text_content if text_content else ["No text could be extracted."]

# Extract text from PPTX
def extract_text_from_pptx(ppt_path):
    prs = Presentation(ppt_path)
    text_content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
    return text_content

#split text into paragraphs
def split_text_into_chunks(text_list, max_length=500):
    chunks = []
    for text in text_list:
        if len(text) > max_length:
            chunks.extend([text[i:i+max_length] for i in range(0, len(text), max_length)])
        else:
            chunks.append(text)
    return chunks

# Function to call OpenRouter API for Q&A generation
def generate_flashcards_from_api(text_list):
    flashcards = []
    text_chunks = split_text_into_chunks(text_list, max_length=500)
    
    for paragraph in text_chunks:
        prompt = f"Generate a question and answer based on this text: {paragraph}"
        
        
        payload = {
            "model": "deepseek/deepseek-chat-v3-0324:free",
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 200
        }

        response = requests.post(api_url, headers=headers, json=payload)
        if response.status_code == 200:
            try:
                response_data = response.json()
                generated_content = response_data["choices"][0]["message"]["content"].strip()

                # Debugging: print the generated content to inspect
                print("Generated Content:", generated_content)

                # Parse the question and answer from the formatted response
                if "**Question:**" in generated_content and "**Answer:**" in generated_content:
                    question_start = generated_content.find("**Question:**") + len("**Question:**") + 1
                    question_end = generated_content.find("**Answer:**")
                    question = generated_content[question_start:question_end].strip()

                    answer_start = generated_content.find("**Answer:**") + len("**Answer:**") + 1
                    answer = generated_content[answer_start:].strip()

                    flashcards.append({"question": question, "answer": answer})

                else:
                    flashcards.append({"question": "Error", "answer": "Unable to parse question and answer."})

            except (KeyError, IndexError) as e:
                flashcards.append({"question": "Error", "answer": f"Failed to parse API response: {str(e)}"})
        else:
            flashcards.append({"question": "Error", "answer": f"Failed with status {response.status_code}"})

    print("Generated flashcards:", flashcards)  # Debugging step
    return flashcards

# Save flashcards to CSV
def save_flashcards_as_csv(flashcards, filename):
    with open(filename, mode="w", newline="", encoding="utf-8") as file:
        writer = csv.writer(file)
        writer.writerow(["Question", "Answer"])  # Header row
        for card in flashcards:
            writer.writerow([card["question"], card["answer"]])

# Save flashcards to JSON
def save_flashcards_as_json(flashcards, filename):
    with open(filename, "w") as f:
        json.dump(flashcards, f, indent=4)

@app.route("/")
def index():
    return render_template("index.html") # changed this from send_from_directory to render_template

@app.route("/flashcards")
def flashcards():
    return render_template("flashcards.html")  # This should point to the flashcards page template


@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400
    
    try:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(file_path)

        # Extract text based on file type
        if file.filename.endswith(".pdf"):
            text_content = extract_text_from_pdf(file_path)
        elif file.filename.endswith(".ppt") or file.filename.endswith(".pptx"):
            text_content = extract_text_from_pptx(file_path)
        else:
            return jsonify({"error": "Unsupported file type"}), 400

        # Generate flashcards using the API
        flashcards = generate_flashcards_from_api(text_content)

        # Save flashcards as CSV & JSON
        csv_path = os.path.join(OUTPUT_FOLDER, "flashcards.csv")
        json_path = os.path.join(OUTPUT_FOLDER, "flashcards.json")
        save_flashcards_as_csv(flashcards, csv_path)
        save_flashcards_as_json(flashcards, json_path)

        return jsonify({
            "csvUrl": "http://localhost:5000/download/csv",
            "jsonUrl": "http://localhost:5000/download/json",
            "flashcards": flashcards
        })
    
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/upload_json", methods=["POST"])
def upload_json_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    if file.filename == "" or not file.filename.endswith(".json"):
        return jsonify({"error": "Please upload a valid JSON file"}), 400

    try:
        # Read the JSON file
        flashcards = json.load(file)

        return jsonify({
            "flashcards": flashcards  # Send back the flashcards data to the frontend
        })
    
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download/<filetype>", methods=["GET"])
def download_file(filetype):
    if filetype == "csv":
        return send_file(os.path.join(OUTPUT_FOLDER, "flashcards.csv"), as_attachment=True)
    elif filetype == "json":
        return send_file(os.path.join(OUTPUT_FOLDER, "flashcards.json"), as_attachment=True)
    else:
        return jsonify({"error": "Invalid file type"}), 400

if __name__ == "__main__":
    app.run(debug=True)